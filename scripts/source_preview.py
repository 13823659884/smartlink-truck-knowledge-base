from __future__ import annotations

import base64
import hashlib
import os
import re
import subprocess
import threading
import zipfile
from pathlib import Path
from typing import Any

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


PREVIEW_WIDTH = 1200
PREVIEW_HEIGHT = 675
RENDER_VERSION = "3"
_PREVIEW_LOCKS: dict[str, threading.Lock] = {}
_PREVIEW_LOCKS_GUARD = threading.Lock()


def locator_number(locator: str, default: int = 1) -> int:
    match = re.search(r"(\d+)", locator or "")
    return max(1, int(match.group(1))) if match else default


def cache_path(
    source_path: Path, locator: str, preview_dir: Path
) -> Path:
    stamp = (
        f"{RENDER_VERSION}|{source_path.resolve()}|"
        f"{source_path.stat().st_mtime_ns}|{locator}"
    )
    digest = hashlib.sha256(stamp.encode("utf-8")).hexdigest()[:24]
    return preview_dir / f"{digest}.png"


def preview_lock(target: Path) -> threading.Lock:
    key = str(target.resolve())
    with _PREVIEW_LOCKS_GUARD:
        lock = _PREVIEW_LOCKS.get(key)
        if lock is None:
            lock = threading.Lock()
            _PREVIEW_LOCKS[key] = lock
        return lock


def chinese_font(size: int) -> ImageFont.FreeTypeFont | ImageFont.ImageFont:
    candidates = [
        Path("C:/Windows/Fonts/msyh.ttc"),
        Path("C:/Windows/Fonts/simhei.ttf"),
        Path("C:/Windows/Fonts/simsun.ttc"),
    ]
    for candidate in candidates:
        if candidate.exists():
            return ImageFont.truetype(str(candidate), size=size)
    return ImageFont.load_default()


def render_pdf_page(source_path: Path, page: int, target: Path) -> None:
    prefix = target.with_suffix("")
    command = [
        "pdftoppm",
        "-f",
        str(page),
        "-l",
        str(page),
        "-singlefile",
        "-png",
        "-r",
        "120",
        str(source_path),
        str(prefix),
    ]
    completed = subprocess.run(
        command,
        capture_output=True,
        timeout=45,
        check=False,
        creationflags=getattr(subprocess, "CREATE_NO_WINDOW", 0),
    )
    if completed.returncode != 0 or not target.exists():
        detail = completed.stderr.decode("utf-8", errors="replace")[:300]
        raise RuntimeError(f"PDF页面渲染失败：{detail}")


def draw_fitted_image(
    canvas: Image.Image,
    image: Image.Image,
    box: tuple[int, int, int, int],
) -> None:
    left, top, right, bottom = box
    width = max(1, right - left)
    height = max(1, bottom - top)
    image = image.convert("RGB")
    image.thumbnail((width, height), Image.Resampling.LANCZOS)
    x = left + (width - image.width) // 2
    y = top + (height - image.height) // 2
    canvas.paste(image, (x, y))


def wrap_text(draw: ImageDraw.ImageDraw, text: str, font: Any, width: int) -> str:
    lines: list[str] = []
    for paragraph in text.splitlines():
        current = ""
        for character in paragraph:
            candidate = current + character
            if draw.textlength(candidate, font=font) > width and current:
                lines.append(current)
                current = character
            else:
                current = candidate
        if current:
            lines.append(current)
        if len(lines) >= 12:
            break
    return "\n".join(lines[:12])


def render_pptx_with_powerpoint(
    source_path: Path, slide_number: int, target: Path
) -> None:
    script = r"""
$ErrorActionPreference = 'Stop'
$application = $null
$presentation = $null
try {
    $application = New-Object -ComObject PowerPoint.Application
    $presentation = $application.Presentations.Open(
        $env:KB_PPT_SOURCE, $true, $true, $false
    )
    $slide = $presentation.Slides.Item([int]$env:KB_PPT_SLIDE)
    $slide.Export($env:KB_PPT_TARGET, 'PNG', 1200, 675)
}
finally {
    if ($presentation -ne $null) { $presentation.Close() }
    if ($application -ne $null) { $application.Quit() }
}
"""
    encoded = base64.b64encode(script.encode("utf-16le")).decode("ascii")
    environment = os.environ.copy()
    environment.update(
        {
            "KB_PPT_SOURCE": str(source_path.resolve()),
            "KB_PPT_SLIDE": str(slide_number),
            "KB_PPT_TARGET": str(target.resolve()),
        }
    )
    completed = subprocess.run(
        [
            "powershell.exe",
            "-NoProfile",
            "-NonInteractive",
            "-STA",
            "-EncodedCommand",
            encoded,
        ],
        capture_output=True,
        timeout=60,
        check=False,
        env=environment,
        creationflags=getattr(subprocess, "CREATE_NO_WINDOW", 0),
    )
    if completed.returncode != 0 or not target.exists():
        detail = completed.stderr.decode("utf-8", errors="replace")[:400]
        raise RuntimeError(f"PowerPoint导出失败：{detail}")


def render_pptx_fallback(
    source_path: Path, slide_number: int, target: Path
) -> None:
    presentation = Presentation(str(source_path))
    if slide_number > len(presentation.slides):
        raise ValueError("幻灯片页码超出范围")
    slide = presentation.slides[slide_number - 1]
    slide_width = float(presentation.slide_width)
    slide_height = float(presentation.slide_height)
    x_scale = PREVIEW_WIDTH / slide_width
    y_scale = PREVIEW_HEIGHT / slide_height
    canvas = Image.new("RGB", (PREVIEW_WIDTH, PREVIEW_HEIGHT), "white")
    draw = ImageDraw.Draw(canvas)
    body_font = chinese_font(22)
    title_font = chinese_font(30)

    for shape in slide.shapes:
        left = max(0, int(shape.left * x_scale))
        top = max(0, int(shape.top * y_scale))
        right = min(PREVIEW_WIDTH, int((shape.left + shape.width) * x_scale))
        bottom = min(PREVIEW_HEIGHT, int((shape.top + shape.height) * y_scale))
        if right <= left or bottom <= top:
            continue
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                from io import BytesIO

                picture = Image.open(BytesIO(shape.image.blob))
                draw_fitted_image(canvas, picture, (left, top, right, bottom))
            except Exception:
                continue
        elif getattr(shape, "has_text_frame", False):
            text = shape.text.strip()
            if not text:
                continue
            font = title_font if top < PREVIEW_HEIGHT * 0.2 else body_font
            box_width = max(1, right - left)
            box_height = max(1, bottom - top)
            text_layer = Image.new(
                "RGBA", (box_width, box_height), (255, 255, 255, 0)
            )
            text_draw = ImageDraw.Draw(text_layer)
            wrapped = wrap_text(text_draw, text, font, max(60, box_width))
            text_draw.multiline_text(
                (0, 0),
                wrapped,
                fill="#1f2937",
                font=font,
                spacing=7,
            )
            canvas.paste(text_layer, (left, top), text_layer)
    draw.rectangle((0, 0, PREVIEW_WIDTH - 1, PREVIEW_HEIGHT - 1), outline="#d1d5db")
    draw.text(
        (PREVIEW_WIDTH - 160, PREVIEW_HEIGHT - 38),
        f"第 {slide_number} 页",
        fill="#6b7280",
        font=chinese_font(18),
    )
    canvas.save(target, format="PNG", optimize=True)


def render_pptx_slide(source_path: Path, slide_number: int, target: Path) -> None:
    try:
        render_pptx_with_powerpoint(source_path, slide_number, target)
    except Exception:
        render_pptx_fallback(source_path, slide_number, target)


def render_docx_image(source_path: Path, target: Path) -> None:
    with zipfile.ZipFile(source_path) as archive:
        media = sorted(
            name
            for name in archive.namelist()
            if name.startswith("word/media/")
            and Path(name).suffix.lower() in {".png", ".jpg", ".jpeg", ".webp", ".bmp"}
        )
        if not media:
            raise ValueError("Word文档中没有可预览图片")
        from io import BytesIO

        image = Image.open(BytesIO(archive.read(media[0])))
    canvas = Image.new("RGB", (PREVIEW_WIDTH, PREVIEW_HEIGHT), "white")
    draw_fitted_image(canvas, image, (20, 20, PREVIEW_WIDTH - 20, PREVIEW_HEIGHT - 20))
    canvas.save(target, format="PNG", optimize=True)


def create_preview(
    source_path: Path, locator: str, preview_dir: Path
) -> Path | None:
    preview_dir.mkdir(parents=True, exist_ok=True)
    suffix = source_path.suffix.lower()
    target = cache_path(source_path, locator, preview_dir)
    if target.exists() and target.stat().st_size > 0:
        return target
    with preview_lock(target):
        if target.exists() and target.stat().st_size > 0:
            return target
        try:
            if suffix == ".pdf":
                render_pdf_page(source_path, locator_number(locator), target)
            elif suffix == ".pptx":
                render_pptx_slide(source_path, locator_number(locator), target)
            elif suffix == ".docx":
                render_docx_image(source_path, target)
            elif suffix in {".png", ".jpg", ".jpeg", ".webp", ".bmp"}:
                image = Image.open(source_path)
                canvas = Image.new("RGB", (PREVIEW_WIDTH, PREVIEW_HEIGHT), "white")
                draw_fitted_image(
                    canvas,
                    image,
                    (10, 10, PREVIEW_WIDTH - 10, PREVIEW_HEIGHT - 10),
                )
                canvas.save(target, format="PNG", optimize=True)
            else:
                return None
        except Exception:
            if target.exists():
                target.unlink()
            raise
    return target if target.exists() else None
