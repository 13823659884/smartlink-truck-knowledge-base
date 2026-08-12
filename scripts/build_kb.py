from __future__ import annotations

import argparse
import csv
import hashlib
import json
import math
import re
import sqlite3
import sys
import traceback
import unicodedata
from collections import Counter, defaultdict
from datetime import datetime, timezone
from pathlib import Path
from typing import Iterable

from docx import Document
from openpyxl import load_workbook
from pypdf import PdfReader
from pptx import Presentation

from ocr_engine import ocr_status, recognize_pdf


BASE_DIR = Path(__file__).resolve().parents[1]
DEFAULT_CONFIG = BASE_DIR / "config.json"
SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".pptx", ".xlsx", ".doc", ".xls"}
INVENTORY_EXTENSIONS = SUPPORTED_EXTENSIONS
EXTRACTION_CACHE_VERSION = 1
OCR_CACHE_VERSION = 1


def utc_now() -> str:
    return datetime.now(timezone.utc).isoformat(timespec="seconds")


def normalize_text(value: object) -> str:
    text = (
        str(value or "")
        .encode("utf-8", errors="replace")
        .decode("utf-8")
        .replace("\x00", " ")
    )
    text = "".join(
        char
        if char in {"\n", "\t"} or unicodedata.category(char) not in {"Cc", "Cs", "Co"}
        else " "
        for char in text
    )
    text = re.sub(r"[ \t\u3000]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def stable_id(prefix: str, value: str) -> str:
    digest = hashlib.sha1(value.encode("utf-8")).hexdigest()[:20]
    return f"{prefix}:{digest}"


def sha256_file(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as stream:
        for block in iter(lambda: stream.read(1024 * 1024), b""):
            digest.update(block)
    return digest.hexdigest()


def parse_version(name: str) -> tuple[str, tuple[int, ...]]:
    match = re.search(r"(?i)(?:^|[^a-z])v(?:ersion)?\s*(\d+(?:[._-]\d+)*)", name)
    if not match:
        return "", ()
    raw = match.group(1).replace("_", ".").replace("-", ".")
    return f"V{raw}", tuple(int(part) for part in raw.split("."))


def parse_date(name: str) -> tuple[str, int]:
    compact = re.search(r"(?<!\d)(20\d{2})(0[1-9]|1[0-2])([0-2]\d|3[01])(?!\d)", name)
    if compact:
        value = "".join(compact.groups())
        return f"{value[:4]}-{value[4:6]}-{value[6:]}", int(value)
    separated = re.search(
        r"(?<!\d)(20\d{2})[-_.年](\d{1,2})[-_.月](\d{1,2})(?:日)?(?!\d)",
        name,
    )
    if separated:
        year, month, day = (int(part) for part in separated.groups())
        value = year * 10000 + month * 100 + day
        return f"{year:04d}-{month:02d}-{day:02d}", value
    return "", 0


def logical_key(relative_path: Path) -> str:
    stem = relative_path.stem.lower()
    stem = re.sub(r"(?i)v(?:ersion)?\s*\d+(?:[._-]\d+)*", "", stem)
    stem = re.sub(r"20\d{2}[-_.年]?\d{1,2}[-_.月]?\d{1,2}日?", "", stem)
    stem = re.sub(r"\((?:新|修订|最终|final|\d+)\)", "", stem, flags=re.I)
    stem = re.sub(r"[\s_\-—–（）()]+", "", stem)
    return f"{relative_path.parent.as_posix().lower()}::{stem}"


def is_cjk(char: str) -> bool:
    return "\u3400" <= char <= "\u9fff"


def search_tokens(text: str) -> list[str]:
    normalized = normalize_text(text).lower()
    tokens: list[str] = []
    seen: set[str] = set()

    def add(token: str) -> None:
        token = token.strip()
        if len(token) < 2 or token in seen:
            return
        seen.add(token)
        tokens.append(token)

    for word in re.findall(r"[a-z0-9]+(?:[-_.][a-z0-9]+)*", normalized):
        add(word)
        for part in re.split(r"[-_.]", word):
            add(part)

    cjk_run = ""
    for char in normalized:
        if is_cjk(char):
            cjk_run += char
        else:
            if cjk_run:
                for width in (2, 3):
                    for index in range(max(0, len(cjk_run) - width + 1)):
                        add(cjk_run[index : index + width])
                cjk_run = ""
    if cjk_run:
        for width in (2, 3):
            for index in range(max(0, len(cjk_run) - width + 1)):
                add(cjk_run[index : index + width])

    return tokens


def vector_counts(text: str, dimensions: int) -> Counter[int]:
    counts: Counter[int] = Counter()
    normalized = re.sub(r"\s+", "", normalize_text(text).lower())
    for width in (2, 3, 4):
        for index in range(max(0, len(normalized) - width + 1)):
            gram = normalized[index : index + width]
            digest = hashlib.blake2b(gram.encode("utf-8"), digest_size=4).digest()
            bucket = int.from_bytes(digest, "big") % dimensions
            counts[bucket] += 1
    for token in re.findall(r"[a-z0-9]+(?:[-_.][a-z0-9]+)*", normalized):
        digest = hashlib.blake2b(token.encode("utf-8"), digest_size=4).digest()
        counts[int.from_bytes(digest, "big") % dimensions] += 2
    return counts


def normalize_vector(counts: Counter[int], idf: list[float]) -> dict[str, float]:
    weighted = {
        index: (1.0 + math.log(count)) * idf[index]
        for index, count in counts.items()
        if count > 0
    }
    norm = math.sqrt(sum(value * value for value in weighted.values())) or 1.0
    return {
        str(index): round(value / norm, 6)
        for index, value in weighted.items()
    }


def split_text(text: str, chunk_size: int, overlap: int) -> list[str]:
    text = normalize_text(text)
    if not text:
        return []
    chunks: list[str] = []
    start = 0
    while start < len(text):
        end = min(len(text), start + chunk_size)
        if end < len(text):
            window = text[start:end]
            candidates = [
                window.rfind(marker)
                for marker in ("\n\n", "\n", "。", "；", "！", "？", ". ")
            ]
            boundary = max(candidates)
            if boundary >= int(chunk_size * 0.55):
                end = start + boundary + 1
        part = text[start:end].strip()
        if part:
            chunks.append(part)
        if end >= len(text):
            break
        start = max(start + 1, end - overlap)
    return chunks


def extract_pdf(path: Path) -> list[tuple[str, str]]:
    reader = PdfReader(str(path))
    units: list[tuple[str, str]] = []
    for index, page in enumerate(reader.pages, start=1):
        text = normalize_text(page.extract_text() or "")
        if text:
            units.append((f"第{index}页", text))
    return units


def extract_docx(path: Path) -> list[tuple[str, str]]:
    document = Document(path)
    units: list[tuple[str, str]] = []
    for index, paragraph in enumerate(document.paragraphs, start=1):
        text = normalize_text(paragraph.text)
        if text:
            units.append((f"段落{index}", text))
    for table_index, table in enumerate(document.tables, start=1):
        for row_index, row in enumerate(table.rows, start=1):
            text = normalize_text(" | ".join(cell.text for cell in row.cells))
            if text:
                units.append((f"表{table_index}第{row_index}行", text))
    return coalesce_units(units)


def extract_pptx(path: Path) -> list[tuple[str, str]]:
    presentation = Presentation(path)
    units: list[tuple[str, str]] = []
    for index, slide in enumerate(presentation.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = normalize_text(shape.text)
                if text:
                    texts.append(text)
        joined = "\n".join(texts)
        if joined:
            units.append((f"第{index}页幻灯片", joined))
    return units


def extract_xlsx(path: Path) -> list[tuple[str, str]]:
    workbook = load_workbook(path, read_only=True, data_only=True)
    units: list[tuple[str, str]] = []
    try:
        for sheet in workbook.worksheets:
            batch: list[str] = []
            start_row = 1
            last_row = 1
            for row_index, row in enumerate(sheet.iter_rows(values_only=True), start=1):
                values = [normalize_text(value) for value in row if value not in (None, "")]
                if not values:
                    continue
                if not batch:
                    start_row = row_index
                batch.append(" | ".join(values))
                last_row = row_index
                if len("\n".join(batch)) >= 2500:
                    units.append(
                        (f"{sheet.title}!第{start_row}-{last_row}行", "\n".join(batch))
                    )
                    batch = []
            if batch:
                units.append(
                    (f"{sheet.title}!第{start_row}-{last_row}行", "\n".join(batch))
                )
    finally:
        workbook.close()
    return units


def extract_legacy_binary(path: Path) -> list[tuple[str, str]]:
    data = path.read_bytes()
    decoded_candidates = [
        data.decode("utf-16le", errors="ignore"),
        data.decode("gb18030", errors="ignore"),
    ]
    fragments: list[str] = []
    seen: set[str] = set()
    pattern = re.compile(r"[\u3400-\u9fffA-Za-z0-9，。；：、（）()/%+\-_.\s]{4,500}")
    total_chars = 0
    for decoded in decoded_candidates:
        for match in pattern.finditer(decoded):
            text = normalize_text(match.group(0))
            if len(text) < 4:
                continue
            meaningful = sum(
                1
                for char in text
                if is_cjk(char) or char.isalnum()
            )
            if meaningful / max(1, len(text)) < 0.72:
                continue
            if text in seen:
                continue
            seen.add(text)
            fragments.append(text)
            total_chars += len(text)
            if total_chars >= 500_000:
                break
        if total_chars >= 500_000:
            break
    units = [
        (f"二进制文本片段{index}", text)
        for index, text in enumerate(fragments, start=1)
    ]
    return coalesce_units(units, target_size=1800)


def coalesce_units(
    units: list[tuple[str, str]], target_size: int = 1200
) -> list[tuple[str, str]]:
    result: list[tuple[str, str]] = []
    locators: list[str] = []
    texts: list[str] = []
    size = 0
    for locator, text in units:
        if texts and size + len(text) > target_size:
            combined_locator = (
                locators[0]
                if len(locators) == 1
                else f"{locators[0]}—{locators[-1]}"
            )
            result.append((combined_locator, "\n".join(texts)))
            locators, texts, size = [], [], 0
        locators.append(locator)
        texts.append(text)
        size += len(text)
    if texts:
        combined_locator = (
            locators[0] if len(locators) == 1 else f"{locators[0]}—{locators[-1]}"
        )
        result.append((combined_locator, "\n".join(texts)))
    return result


def extract_units(path: Path) -> list[tuple[str, str]]:
    extension = path.suffix.lower()
    if extension == ".pdf":
        return extract_pdf(path)
    if extension == ".docx":
        return extract_docx(path)
    if extension == ".pptx":
        return extract_pptx(path)
    if extension == ".xlsx":
        return extract_xlsx(path)
    if extension in {".doc", ".xls"}:
        return extract_legacy_binary(path)
    return []


def cached_extract_units(
    path: Path,
    file_sha256: str,
    cache_dir: Path,
    ocr_config: dict[str, object],
) -> tuple[list[tuple[str, str]], bool, bool, bool, int]:
    """Reuse expensive document extraction when file content is unchanged."""
    cache_dir.mkdir(parents=True, exist_ok=True)
    cache_path = cache_dir / f"v{EXTRACTION_CACHE_VERSION}-{file_sha256}.json"
    cached_payload: dict[str, object] | None = None
    raw_cache_hit = False
    if cache_path.exists():
        try:
            payload = json.loads(cache_path.read_text(encoding="utf-8"))
            if (
                payload.get("version") == EXTRACTION_CACHE_VERSION
                and payload.get("extension") == path.suffix.lower()
            ):
                units = [
                    (str(item[0]), str(item[1]))
                    for item in payload.get("units", [])
                    if isinstance(item, list) and len(item) == 2
                ]
                cached_payload = payload
                raw_cache_hit = True
                ocr_meta = payload.get("ocr", {})
                if not isinstance(ocr_meta, dict):
                    ocr_meta = {}
                ocr_used = bool(ocr_meta.get("used"))
                ocr_processed = (
                    int(ocr_meta.get("version", 0) or 0) == OCR_CACHE_VERSION
                )
                ocr_pages = int(ocr_meta.get("pages", 0) or 0)
                ocr_enabled = bool(ocr_config.get("enabled", False))
                ocr_current = int(ocr_meta.get("version", 0) or 0) == OCR_CACHE_VERSION
                if (
                    units
                    or path.suffix.lower() != ".pdf"
                    or not ocr_enabled
                    or ocr_current
                ):
                    return units, True, ocr_used, ocr_processed, ocr_pages
        except (OSError, ValueError, TypeError, json.JSONDecodeError):
            pass

    units = [] if raw_cache_hit else extract_units(path)
    ocr_meta: dict[str, object] = {}
    ocr_attempted = False
    if (
        not units
        and path.suffix.lower() == ".pdf"
        and bool(ocr_config.get("enabled", False))
        and bool(ocr_config.get("runtime_available", False))
    ):
        ocr_attempted = True
        result = recognize_pdf(
            path,
            dpi=int(ocr_config.get("render_dpi", 160)),
            min_confidence=float(ocr_config.get("min_confidence", 0.45)),
            max_pages=int(ocr_config.get("max_pages_per_document", 200)),
        )
        units = [
            (str(locator), str(text))
            for locator, text in result.get("units", [])
        ]
        ocr_meta = {
            "version": OCR_CACHE_VERSION,
            "used": bool(units),
            "pages": int(result.get("pages", 0)),
            "recognized_pages": int(result.get("recognized_pages", 0)),
            "confidence": float(result.get("confidence", 0.0)),
            "elapsed_ms": float(result.get("elapsed_ms", 0.0)),
            "engine": str(result.get("engine", "RapidOCR ONNX")),
        }
    elif cached_payload and isinstance(cached_payload.get("ocr"), dict):
        ocr_meta = dict(cached_payload["ocr"])

    cache_path.write_text(
        json.dumps(
            {
                "version": EXTRACTION_CACHE_VERSION,
                "extension": path.suffix.lower(),
                "units": units,
                "ocr": ocr_meta,
            },
            ensure_ascii=False,
            separators=(",", ":"),
        ),
        encoding="utf-8",
    )
    return (
        units,
        raw_cache_hit and not ocr_attempted,
        bool(ocr_meta.get("used")),
        int(ocr_meta.get("version", 0) or 0) == OCR_CACHE_VERSION,
        int(ocr_meta.get("pages", 0) or 0),
    )


def detect_scene(relative_path: Path, scene_mapping: dict[str, str]) -> str:
    for part in relative_path.parts:
        if part in scene_mapping:
            return scene_mapping[part]
    return "其他"


def detect_tags(text: str, configured_series: list[str]) -> tuple[list[str], list[str]]:
    upper = text.upper()
    vehicle_tags = [
        series
        for series in configured_series
        if series.upper() in upper
    ]
    energy_tags: list[str] = []
    if re.search(r"新能源|纯电|电车|换电|燃电|动力电池|高压", text, re.I):
        energy_tags.append("新能源")
    if re.search(r"传统车|燃油|柴油|发动机", text, re.I):
        energy_tags.append("传统")
    return sorted(set(vehicle_tags)), sorted(set(energy_tags))


def create_schema(connection: sqlite3.Connection, schema_path: Path) -> None:
    connection.executescript(schema_path.read_text(encoding="utf-8"))


def snapshot_runtime_data(database_path: Path) -> dict[str, list[dict[str, object]]]:
    preserved: dict[str, list[dict[str, object]]] = {
        "conversations": [],
        "messages": [],
        "feedback": [],
        "query_intents": [],
        "vin_records": [],
    }
    if not database_path.exists():
        return preserved
    connection = sqlite3.connect(database_path)
    connection.row_factory = sqlite3.Row
    try:
        for table in preserved:
            try:
                preserved[table] = [
                    dict(row)
                    for row in connection.execute(f"SELECT * FROM {table}")
                ]
            except sqlite3.OperationalError:
                preserved[table] = []
    finally:
        connection.close()
    return preserved


def restore_runtime_data(
    connection: sqlite3.Connection,
    preserved: dict[str, list[dict[str, object]]],
) -> dict[str, int]:
    restored: dict[str, int] = {}
    for table in (
        "conversations",
        "messages",
        "feedback",
        "query_intents",
        "vin_records",
    ):
        rows = preserved.get(table, [])
        if not rows:
            restored[table] = 0
            continue
        columns = list(rows[0].keys())
        column_sql = ", ".join(columns)
        placeholders = ", ".join("?" for _ in columns)
        connection.executemany(
            f"INSERT INTO {table} ({column_sql}) VALUES ({placeholders})",
            [tuple(row.get(column) for column in columns) for row in rows],
        )
        restored[table] = len(rows)
    return restored


def import_graph(
    connection: sqlite3.Connection, graph_output: Path, triples_path: Path
) -> dict[str, int]:
    node_path = graph_output / "nodes.csv"
    relationship_path = graph_output / "relationships.csv"
    nodes: dict[str, dict[str, str]] = {}

    with node_path.open("r", encoding="utf-8-sig", newline="") as stream:
        for row in csv.DictReader(stream):
            nodes[row["id"]] = row
            connection.execute(
                """
                INSERT OR REPLACE INTO entities
                (id, label, name, code, description, source_path, source_locator,
                 confidence, review_status)
                VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?)
                """,
                (
                    row["id"],
                    row["label"],
                    row["name"],
                    row["code"],
                    row["description"],
                    row["source_path"],
                    row["source_locator"],
                    row["confidence"],
                    row["review_status"],
                ),
            )

    exported: list[dict[str, str]] = []
    with relationship_path.open("r", encoding="utf-8-sig", newline="") as stream:
        for row in csv.DictReader(stream):
            subject = nodes.get(row["start_id"], {})
            obj = nodes.get(row["end_id"], {})
            if not subject or not obj:
                continue
            record = {
                "id": row["id"],
                "subject_id": row["start_id"],
                "subject": subject.get("name", ""),
                "subject_label": subject.get("label", ""),
                "predicate": row["type"],
                "object_id": row["end_id"],
                "object": obj.get("name", ""),
                "object_label": obj.get("label", ""),
                "source_path": row["source_path"],
                "source_locator": row["source_locator"],
                "confidence": row["confidence"],
                "review_status": row["review_status"],
            }
            exported.append(record)
            connection.execute(
                """
                INSERT OR REPLACE INTO triples
                (id, subject_id, subject, subject_label, predicate, object_id,
                 object, object_label, source_path, source_locator, confidence,
                 review_status)
                VALUES (:id, :subject_id, :subject, :subject_label, :predicate,
                        :object_id, :object, :object_label, :source_path,
                        :source_locator, :confidence, :review_status)
                """,
                record,
            )
            graph_search_terms = " ".join(
                search_tokens(
                    " ".join(
                        [
                            record["subject"],
                            record["subject_label"],
                            record["predicate"],
                            record["object"],
                            record["object_label"],
                        ]
                    )
                )
            )
            connection.execute(
                "DELETE FROM triples_fts WHERE triple_id = ?", (record["id"],)
            )
            connection.execute(
                "INSERT INTO triples_fts(triple_id, search_terms) VALUES (?, ?)",
                (record["id"], graph_search_terms),
            )

    triples_path.parent.mkdir(parents=True, exist_ok=True)
    if exported:
        with triples_path.open("w", encoding="utf-8-sig", newline="") as stream:
            writer = csv.DictWriter(stream, fieldnames=list(exported[0].keys()))
            writer.writeheader()
            writer.writerows(exported)
    return {"entities": len(nodes), "triples": len(exported)}


def build(config_path: Path) -> dict[str, object]:
    started_at = utc_now()
    config = json.loads(config_path.read_text(encoding="utf-8"))
    source_root = (BASE_DIR / config["source_root"]).resolve()
    graph_output = (BASE_DIR / config["graph_output"]).resolve()
    output_dir = (BASE_DIR / config["output_dir"]).resolve()
    output_dir.mkdir(parents=True, exist_ok=True)
    database_path = output_dir / "knowledge_base.db"
    report_path = output_dir / "build_report.json"
    unparsed_path = output_dir / "unparsed_files.csv"
    triples_path = output_dir / "triples.csv"
    extraction_cache_dir = output_dir / "extraction_cache"
    ocr_config = dict(config.get("ocr", {}))
    ocr_runtime = ocr_status()
    ocr_config["runtime_available"] = bool(ocr_runtime.get("available"))

    preserved_runtime = snapshot_runtime_data(database_path)
    if database_path.exists():
        database_path.unlink()
    connection = sqlite3.connect(database_path)
    connection.row_factory = sqlite3.Row
    create_schema(connection, BASE_DIR / "schema.sql")
    run_id = connection.execute(
        "INSERT INTO build_runs(started_at, status) VALUES (?, ?)",
        (started_at, "running"),
    ).lastrowid

    for scene in sorted(set(config["scenes"].values())):
        connection.execute(
            "INSERT INTO knowledge_bases(id, name, vehicle_series, scene) VALUES (?, ?, ?, ?)",
            (f"ALL|{scene}", f"全车系-{scene}", "", scene),
        )
        for series in config["vehicle_series"]:
            connection.execute(
                "INSERT INTO knowledge_bases(id, name, vehicle_series, scene) VALUES (?, ?, ?, ?)",
                (f"{series}|{scene}", f"{series}-{scene}", series, scene),
            )

    paths = sorted(
        path
        for path in source_root.rglob("*")
        if path.is_file()
        and path.suffix.lower() in INVENTORY_EXTENSIONS
        and not path.name.startswith("~$")
    )
    inventory: list[dict[str, object]] = []
    groups: defaultdict[str, list[dict[str, object]]] = defaultdict(list)
    for path in paths:
        relative_path = path.relative_to(source_root)
        version, version_tuple = parse_version(path.name)
        effective_date, date_number = parse_date(path.name)
        stat = path.stat()
        record: dict[str, object] = {
            "path": path,
            "relative_path": relative_path,
            "logical_key": logical_key(relative_path),
            "version": version,
            "version_tuple": version_tuple,
            "effective_date": effective_date,
            "date_number": date_number,
            "mtime": int(stat.st_mtime),
            "size": stat.st_size,
            "sha256": sha256_file(path),
        }
        inventory.append(record)
        groups[str(record["logical_key"])].append(record)

    active_paths: set[Path] = set()
    for records in groups.values():
        selected = max(
            records,
            key=lambda item: (
                item["version_tuple"],
                item["date_number"],
                item["mtime"],
            ),
        )
        active_paths.add(selected["path"])

    chunk_records: list[dict[str, object]] = []
    unparsed: list[dict[str, str]] = []
    extraction_status: Counter[str] = Counter()
    scene_counts: Counter[str] = Counter()
    extension_counts: Counter[str] = Counter()
    extraction_cache_hits = 0
    extraction_cache_misses = 0
    ocr_documents = 0
    ocr_recognized_documents = 0
    ocr_pages = 0

    for item in inventory:
        path = item["path"]
        relative_path = item["relative_path"]
        assert isinstance(path, Path)
        assert isinstance(relative_path, Path)
        enabled = path in active_paths
        extension = path.suffix.lower()
        scene = detect_scene(relative_path, config["scenes"])
        path_vehicle_tags, path_energy_tags = detect_tags(
            relative_path.as_posix(), config["vehicle_series"]
        )
        document_id = stable_id("document", relative_path.as_posix())
        status = "inactive_version" if not enabled else "pending"
        error_message = ""
        document_chunks: list[dict[str, object]] = []

        if enabled and extension in SUPPORTED_EXTENSIONS:
            try:
                (
                    units,
                    cache_hit,
                    ocr_used,
                    ocr_processed,
                    document_ocr_pages,
                ) = cached_extract_units(
                    path,
                    str(item["sha256"]),
                    extraction_cache_dir,
                    ocr_config,
                )
                if cache_hit:
                    extraction_cache_hits += 1
                else:
                    extraction_cache_misses += 1
                if ocr_processed:
                    ocr_documents += 1
                    ocr_pages += document_ocr_pages
                if ocr_used:
                    ocr_recognized_documents += 1
                if not units:
                    status = "image_only" if ocr_processed else "needs_ocr"
                    unparsed.append(
                        {
                            "relative_path": relative_path.as_posix(),
                            "reason": (
                                "已完成OCR，但未识别到足够清晰的可检索文字"
                                if ocr_processed
                                else "未提取到可检索文本，可能是扫描件、图纸或图片型文档"
                            ),
                        }
                    )
                else:
                    preview_text = "\n".join(text for _, text in units)[:20000]
                    text_vehicle_tags, text_energy_tags = detect_tags(
                        preview_text, config["vehicle_series"]
                    )
                    vehicle_tags = sorted(
                        set(path_vehicle_tags) | set(text_vehicle_tags)
                    )
                    energy_tags = sorted(
                        set(path_energy_tags) | set(text_energy_tags)
                    )
                    ordinal = 0
                    for locator, text in units:
                        for part in split_text(
                            text,
                            int(config["chunk_size"]),
                            int(config["chunk_overlap"]),
                        ):
                            if len(part) < (4 if ocr_used else 20):
                                continue
                            ordinal += 1
                            chunk_id = stable_id(
                                "chunk", f"{relative_path.as_posix()}|{ordinal}|{part}"
                            )
                            tokens = search_tokens(
                                f"{relative_path.as_posix()} {part}"
                            )
                            document_chunks.append(
                                {
                                    "id": chunk_id,
                                    "document_id": document_id,
                                    "ordinal": ordinal,
                                    "source_locator": locator,
                                    "content": part,
                                    "search_terms": " ".join(tokens[:1800]),
                                    "vehicle_tags": ",".join(vehicle_tags),
                                    "scene": scene,
                                    "energy_tags": ",".join(energy_tags),
                                    "token_count": len(tokens),
                                }
                            )
                    status = (
                        "parsed_legacy"
                        if document_chunks and extension in {".doc", ".xls"}
                        else "parsed_ocr"
                        if document_chunks and ocr_used
                        else "parsed"
                        if document_chunks
                        else "image_only"
                        if ocr_processed
                        else "needs_ocr"
                    )
                    if status == "image_only":
                        unparsed.append(
                            {
                                "relative_path": relative_path.as_posix(),
                                "reason": "OCR识别文字过少，未形成有效检索分块",
                            }
                        )
            except Exception as exc:
                status = "error"
                error_message = f"{type(exc).__name__}: {exc}"
                unparsed.append(
                    {
                        "relative_path": relative_path.as_posix(),
                        "reason": error_message,
                    }
                )
        connection.execute(
            """
            INSERT INTO documents
            (id, relative_path, file_name, extension, sha256, size_bytes,
             modified_at, logical_key, version, effective_date, scene,
             vehicle_tags, energy_tags, status, enabled, chunk_count,
             error_message)
            VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?)
            """,
            (
                document_id,
                relative_path.as_posix(),
                path.name,
                extension,
                item["sha256"],
                item["size"],
                datetime.fromtimestamp(
                    item["mtime"], timezone.utc
                ).isoformat(timespec="seconds"),
                item["logical_key"],
                item["version"],
                item["effective_date"],
                scene,
                ",".join(path_vehicle_tags),
                ",".join(path_energy_tags),
                status,
                int(enabled),
                len(document_chunks),
                error_message,
            ),
        )
        chunk_records.extend(document_chunks)
        extraction_status[status] += 1
        scene_counts[scene] += 1
        extension_counts[extension] += 1

    dimensions = int(config["vector_dimensions"])
    chunk_counts: dict[str, Counter[int]] = {}
    document_frequency: Counter[int] = Counter()
    for chunk in chunk_records:
        counts = vector_counts(str(chunk["content"]), dimensions)
        chunk_counts[str(chunk["id"])] = counts
        document_frequency.update(counts.keys())
    total_chunks = max(1, len(chunk_records))
    idf = [
        math.log((total_chunks + 1) / (document_frequency[index] + 1)) + 1.0
        for index in range(dimensions)
    ]
    connection.execute(
        "INSERT INTO vector_meta(key, value_json) VALUES (?, ?)",
        ("idf", json.dumps(idf, separators=(",", ":"))),
    )
    connection.execute(
        "INSERT INTO vector_meta(key, value_json) VALUES (?, ?)",
        ("dimensions", json.dumps(dimensions)),
    )

    for chunk in chunk_records:
        connection.execute(
            """
            INSERT INTO chunks
            (id, document_id, ordinal, source_locator, content, search_terms,
             vehicle_tags, scene, energy_tags, token_count)
            VALUES (:id, :document_id, :ordinal, :source_locator, :content,
                    :search_terms, :vehicle_tags, :scene, :energy_tags,
                    :token_count)
            """,
            chunk,
        )
        connection.execute(
            "INSERT INTO chunks_fts(chunk_id, search_terms, content) VALUES (?, ?, ?)",
            (chunk["id"], chunk["search_terms"], chunk["content"]),
        )
        vector = normalize_vector(chunk_counts[str(chunk["id"])], idf)
        connection.execute(
            "INSERT INTO chunk_vectors(chunk_id, vector_json) VALUES (?, ?)",
            (chunk["id"], json.dumps(vector, separators=(",", ":"))),
        )

    graph_stats = import_graph(connection, graph_output, triples_path)
    restored_runtime = restore_runtime_data(connection, preserved_runtime)

    with unparsed_path.open("w", encoding="utf-8-sig", newline="") as stream:
        writer = csv.DictWriter(stream, fieldnames=["relative_path", "reason"])
        writer.writeheader()
        writer.writerows(unparsed)

    report: dict[str, object] = {
        "valid": True,
        "started_at": started_at,
        "finished_at": utc_now(),
        "source_root": str(source_root),
        "database": str(database_path),
        "documents": len(inventory),
        "active_documents": len(active_paths),
        "chunks": len(chunk_records),
        "entities": graph_stats["entities"],
        "triples": graph_stats["triples"],
        "knowledge_bases": connection.execute(
            "SELECT COUNT(*) FROM knowledge_bases"
        ).fetchone()[0],
        "extraction_status": dict(sorted(extraction_status.items())),
        "documents_by_scene": dict(sorted(scene_counts.items())),
        "documents_by_extension": dict(sorted(extension_counts.items())),
        "unparsed_files": len(unparsed),
        "vector_dimensions": dimensions,
        "extraction_cache": {
            "version": EXTRACTION_CACHE_VERSION,
            "hits": extraction_cache_hits,
            "misses": extraction_cache_misses,
        },
        "ocr": {
            **ocr_runtime,
            "enabled": bool(ocr_config.get("enabled", False)),
            "documents": ocr_documents,
            "recognized_documents": ocr_recognized_documents,
            "pages": ocr_pages,
            "cache_version": OCR_CACHE_VERSION,
        },
        "restored_runtime": restored_runtime,
    }
    connection.execute(
        """
        UPDATE build_runs
        SET finished_at = ?, status = ?, report_json = ?
        WHERE id = ?
        """,
        (
            report["finished_at"],
            "succeeded",
            json.dumps(report, ensure_ascii=False),
            run_id,
        ),
    )
    connection.commit()
    connection.close()
    report_path.write_text(
        json.dumps(report, ensure_ascii=False, indent=2),
        encoding="utf-8",
    )
    return report


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("--config", type=Path, default=DEFAULT_CONFIG)
    parser.add_argument(
        "--skip-qdrant",
        action="store_true",
        help="只重建 SQLite，不同步本地 Qdrant 语义索引",
    )
    args = parser.parse_args()
    try:
        report = build(args.config.resolve())
        active_config = json.loads(args.config.read_text(encoding="utf-8"))
        if (
            not args.skip_qdrant
            and active_config.get("retrieval", {}).get("backend")
            == "qdrant_hybrid"
        ):
            from qdrant_store import build_qdrant_index

            report["qdrant"] = build_qdrant_index(
                Path(str(report["database"]))
            )
    except Exception:
        traceback.print_exc()
        return 1
    print(json.dumps(report, ensure_ascii=False, indent=2))
    return 0


if __name__ == "__main__":
    sys.exit(main())
